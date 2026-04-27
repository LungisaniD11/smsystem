#!/usr/bin/env python3
"""
FutureTech Student Management System - User Manual PowerPoint Generator

Usage: python create_pptx.py

Requirements: pip install python-pptx Pillow

This script creates a professional 14-slide PowerPoint presentation
with user manual and step-by-step guides for the FutureTech SMS.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Create the FutureTech UserManual PowerPoint presentation."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    TEAL = RGBColor(15, 118, 110)
    DARK_GRAY = RGBColor(50, 50, 50)
    LIGHT_GRAY = RGBColor(80, 80, 80)
    
    def add_title_slide(title, subtitle):
        """Add a title slide."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_GRAY
        
        return slide
    
    def add_content_slide(title, content):
        """Add a content slide with title and bullet points."""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = title
        p = title_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TEAL
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = content
        
        # Format paragraphs
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = DARK_GRAY
            paragraph.space_before = Pt(2)
            paragraph.space_after = Pt(2)
        
        return slide
    
    # Slide 1: Title
    print("Creating slide 1: Title...")
    add_title_slide("FutureTech Academy", "Student Management System\nUser Manual v1.0")
    
    # Slide 2: Overview
    print("Creating slide 2: Overview...")
    add_content_slide("System Overview",
        "Secure cloud-based platform for managing student records, profiles, and enrollment data.\n\n"
        "Key Features:\n"
        "- Comprehensive Dashboard\n"
        "- Google OAuth 2.0 Authentication\n"
        "- Complete CRUD Operations\n"
        "- Profile Image Management\n"
        "- Search and Pagination\n"
        "- Enterprise Azure Infrastructure")
    
    # Slide 3: Getting Started
    print("Creating slide 3: Getting Started...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "Getting Started: Quick Access"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEAL
    
    # Add QR code if exists
    qr_path = "futuretech_qr.png"
    if os.path.exists(qr_path):
        slide.shapes.add_picture(qr_path, Inches(3.5), Inches(1.5), width=Inches(3))
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    tf = content_box.text_frame
    tf.text = "Scan QR code or visit:\nhttps://futuretechsmsld-fnfgf9evhpghgzbv.southafricanorth-01.azurewebsites.net/"
    tf.paragraphs[0].font.size = Pt(14)
    
    # Slide 4: Authentication
    print("Creating slide 4: Authentication...")
    add_content_slide("Step 1: User Authentication",
        "Login Steps:\n\n"
        "1. Open application via QR code or URL\n"
        "2. Click Login with Google\n"
        "3. Select your Google account\n"
        "4. Grant permissions\n"
        "5. Redirected to dashboard if authorized\n\n"
        "Note: Only registered admin emails can access.")
    
    # Slide 5: Dashboard
    print("Creating slide 5: Dashboard...")
    add_content_slide("Step 2: Dashboard Overview",
        "Dashboard Features:\n\n"
        "- Active Students Count\n"
        "- Inactive Students Count\n"
        "- Students with Profile Images\n"
        "- New Students This Week\n\n"
        "Functionality:\n"
        "- Quick access to recent students\n"
        "- Search by name or student ID\n"
        "- Navigation menu for all functions")
    
    # Slide 6: Adding Students
    print("Creating slide 6: Adding Students...")
    add_content_slide("Step 3: Adding a New Student",
        "Create Student Record:\n\n"
        "1. Click Add New Student from dashboard\n"
        "2. Enter: First Name, Last Name, Email, Mobile\n"
        "3. Select Enrollment Status\n"
        "4. Upload profile image (JPEG/PNG, max 5MB, optional)\n"
        "5. Click Save\n\n"
        "Formats: JPEG and PNG only\n"
        "File signature verification prevents spoofing")
    
    # Slide 7: Searching
    print("Creating slide 7: Searching...")
    add_content_slide("Step 4: Searching & Viewing Students",
        "Search and View Records:\n\n"
        "1. Use search bar by Name or Student ID\n"
        "2. Results paginated (10 per page)\n"
        "3. Click student name for full details\n"
        "4. Profile images shown with secure links\n\n"
        "Security: Image links expire after 15 minutes\n"
        "Prevents unauthorized access to student profiles")
    
    # Slide 8: Editing
    print("Creating slide 8: Editing...")
    add_content_slide("Step 5: Editing Student Records",
        "Update Student Information:\n\n"
        "1. Find student via search or browse\n"
        "2. Click Edit button\n"
        "3. Update any field (names, email, phone, status)\n"
        "4. Replace profile image if needed\n"
        "5. Click Save Changes\n\n"
        "Tip: Leave image field empty to keep existing image")
    
    # Slide 9: Managing Status
    print("Creating slide 9: Managing Status...")
    add_content_slide("Step 6: Managing Student Status",
        "Two Options Available:\n\n"
        "SOFT DELETE (Recommended):\n"
        "- Marks student as inactive\n"
        "- Record remains in system\n"
        "- Data is preserved\n"
        "- Can be reactivated\n\n"
        "PERMANENT DELETE (Use with caution):\n"
        "- Completely removes record\n"
        "- Cannot be recovered\n"
        "- Requires confirmation")
    
    # Slide 10: Image Management
    print("Creating slide 10: Image Management...")
    add_content_slide("Step 7: Profile Image Management",
        "Image Upload Guidelines:\n\n"
        "Supported: JPEG (.jpg, .jpeg) and PNG (.png)\n"
        "Size Limit: Maximum 5 MB per image\n"
        "Validation: File signature verification\n"
        "Processing: Images resized for efficiency\n\n"
        "Upload Process:\n"
        "- Click Choose Image or drag-and-drop\n"
        "- System validates format and integrity\n"
        "- Image uploaded to secure cloud storage")
    
    # Slide 11: Security
    print("Creating slide 11: Security...")
    add_content_slide("Security Features",
        "Your Data is Protected By:\n\n"
        "- Google OAuth 2.0 (Industry standard)\n"
        "- Admin Whitelist (Approved users only)\n"
        "- HTTPS Encryption (Transit protection)\n"
        "- Anti-CSRF Tokens (Attack protection)\n"
        "- Security Headers (HSTS, CSP)\n"
        "- SAS Links (15-minute image expiry)\n"
        "- Azure Cosmos DB (Enterprise protection)\n"
        "- Azure Blob Storage (Secure media)")
    
    # Slide 12: Best Practices
    print("Creating slide 12: Best Practices...")
    add_content_slide("Best Practices & Tips",
        "Usage Guidelines:\n\n"
        "- Security: Never share login credentials\n"
        "- Data Quality: Enter accurate information\n"
        "- Backups: System auto-backs up\n"
        "- Soft Delete First: Use soft delete by default\n"
        "- Verify Updates: Confirm changes saved\n"
        "- Support: Contact IT for technical issues\n"
        "- Audit Trail: Use your own account only")
    
    # Slide 13: Troubleshooting
    print("Creating slide 13: Troubleshooting...")
    add_content_slide("Troubleshooting Guide",
        "Common Issues:\n\n"
        "Cannot Login:\n"
        "- Verify account registered with admin\n"
        "- Check email domain is whitelisted\n"
        "- Clear browser cookies\n\n"
        "Image Upload Failed:\n"
        "- Check JPEG or PNG format\n"
        "- Verify file size (max 5MB)\n"
        "- Ensure file not corrupted\n\n"
        "Record Not Found:\n"
        "- Try different search terms\n"
        "- Check if marked as inactive")
    
    # Slide 14: Summary
    print("Creating slide 14: Summary...")
    add_content_slide("You're Ready to Use FutureTech SMS!",
        "Key Capabilities:\n\n"
        "- Secure login with Google OAuth\n"
        "- Full student record management (CRUD)\n"
        "- Profile images with secure access\n"
        "- Search and pagination support\n"
        "- Enterprise-grade cloud infrastructure\n\n"
        "Need Help?\n"
        "Contact your system administrator\n\n"
        "FutureTech Academy Student Management System | v1.0")
    
    # Save presentation
    output_path = "FutureTech-UserManual.pptx"
    print(f"\nSaving presentation to: {output_path}")
    prs.save(output_path)
    
    print(f"[SUCCESS] PowerPoint presentation created!")
    print(f"Location: {os.path.abspath(output_path)}")
    print(f"Size: {os.path.getsize(output_path) / 1024:.2f} KB")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("[ERROR] Required libraries not found.")
        print("Install with: pip install python-pptx Pillow")
    except Exception as e:
        print(f"[ERROR] {e}")
